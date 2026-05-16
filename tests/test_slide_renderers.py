import io
import zipfile

import pytest
from docx import Document
from pptx import Presentation

from src.schemas import GenerateSlidesRequestSchema, SlideKind
from src.services import (
    SlideDocxService,
    SlideMarkdownService,
    SlideParserService,
    SlidePptxService,
)

SAMPLE_TEXT: str = (
    "Aula 1: Lógica Digital\n"
    "Introdução\n"
    "Texto introdutório com informação relevante para os alunos.\n"
    "\n"
    "Notas: Lembrar do exemplo do LED.\n"
    "\n"
    "Portas Lógicas\n"
    "Porta NOT: inverte a entrada.\n"
    "\n"
    "Porta AND: precisa de todas as entradas em 1.\n"
    "\n"
    "Notas: Mostrar tabela verdade.\n"
    "\n"
    "Conclusão da Unidade\n"
    "Encerramento da unidade com síntese dos pontos.\n"
)


@pytest.mark.asyncio
async def test_parser_attaches_notes_to_paragraph_and_bullet() -> None:
    """Notas should attach to the paragraph or bullet immediately above."""
    parser: SlideParserService = SlideParserService()
    deck = await parser.parse(SAMPLE_TEXT)
    paragraph_slide = next(
        s for s in deck.slides if s.kind == SlideKind.CONTENT
    )
    bullet_slide = next(
        s for s in deck.slides if s.kind == SlideKind.BULLETS
    )
    assert paragraph_slide.notes == "Lembrar do exemplo do LED."
    assert bullet_slide.notes == "Mostrar tabela verdade."


@pytest.mark.asyncio
async def test_pptx_service_produces_valid_pptx() -> None:
    """The PPTX renderer should emit a valid Office Open XML file."""
    parser: SlideParserService = SlideParserService()
    deck = await parser.parse(SAMPLE_TEXT)
    service: SlidePptxService = SlidePptxService()
    payload: bytes = await service.render(deck=deck)
    assert payload[:2] == b"PK"
    presentation = Presentation(io.BytesIO(payload))
    assert len(presentation.slides) == len(deck.slides)
    notes = [
        slide.notes_slide.notes_text_frame.text
        for slide in presentation.slides
        if slide.has_notes_slide
    ]
    assert any("LED" in note for note in notes)


@pytest.mark.asyncio
async def test_docx_handout_has_one_break_per_slide() -> None:
    """The DOCX handout should hold a notes section for each slide."""
    parser: SlideParserService = SlideParserService()
    deck = await parser.parse(SAMPLE_TEXT)
    service: SlideDocxService = SlideDocxService()
    payload: bytes = await service.render(deck=deck)
    assert payload[:2] == b"PK"
    document = Document(io.BytesIO(payload))
    headers: int = sum(
        1
        for paragraph in document.paragraphs
        if paragraph.text == "Notas do apresentador"
    )
    assert headers == len(deck.slides)


@pytest.mark.asyncio
async def test_markdown_uses_revealjs_separators() -> None:
    """Reveal.js Markdown should separate slides with `---` lines."""
    parser: SlideParserService = SlideParserService()
    deck = await parser.parse(SAMPLE_TEXT)
    service: SlideMarkdownService = SlideMarkdownService()
    output: str = await service.render(deck=deck)
    assert output.count("\n---\n") == len(deck.slides) - 1
    assert "Note:" in output
    assert "## Conclusão da Unidade" in output


@pytest.mark.asyncio
async def test_pptx_uses_widescreen_dimensions() -> None:
    """The PPTX should be 16:9 widescreen (13.333 x 7.5 inches)."""
    parser: SlideParserService = SlideParserService()
    deck = await parser.parse(SAMPLE_TEXT)
    service: SlidePptxService = SlidePptxService()
    payload: bytes = await service.render(deck=deck)
    presentation = Presentation(io.BytesIO(payload))
    assert round(presentation.slide_width.inches, 2) == 13.33
    assert round(presentation.slide_height.inches, 2) == 7.5


@pytest.mark.asyncio
async def test_pptx_zip_contains_expected_parts() -> None:
    """A valid PPTX file is a ZIP with [Content_Types].xml at the root."""
    parser: SlideParserService = SlideParserService()
    deck = await parser.parse(SAMPLE_TEXT)
    service: SlidePptxService = SlidePptxService()
    payload: bytes = await service.render(deck=deck)
    with zipfile.ZipFile(io.BytesIO(payload)) as archive:
        names = archive.namelist()
    assert "[Content_Types].xml" in names
