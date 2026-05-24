import asyncio
import io

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from src.schemas import SlideDeckSchema, SlideKind, SlideSchema

_HERO_LAYOUT_INDEX: int = 6  # blank layout
_SLIDE_WIDTH_INCHES: float = 13.333
_SLIDE_HEIGHT_INCHES: float = 7.5
_BODY_FONT: str = "Calibri"
_TITLE_FONT: str = "Calibri"


class _ThemeColors:
	"""Container for the colors applied across a presentation theme."""

	def __init__(
		self,
		background: RGBColor,
		text: RGBColor,
		muted: RGBColor,
		accent: RGBColor,
		accent_soft: RGBColor,
		hero_background: RGBColor,
		hero_text: RGBColor,
		bullet_background: RGBColor,
	) -> None:
		"""Initialize the bundle of colors used by the renderer.

		Args:
		    background: Default slide background color.
		    text: Primary text color.
		    muted: Color used for footers and eyebrows.
		    accent: Brand color used in titles and accents.
		    accent_soft: Subtle accent variant for backgrounds.
		    hero_background: Background of lesson title slides.
		    hero_text: Text color on hero slides.
		    bullet_background: Background of bullet item boxes.
		"""
		self.background: RGBColor = background
		self.text: RGBColor = text
		self.muted: RGBColor = muted
		self.accent: RGBColor = accent
		self.accent_soft: RGBColor = accent_soft
		self.hero_background: RGBColor = hero_background
		self.hero_text: RGBColor = hero_text
		self.bullet_background: RGBColor = bullet_background


_THEMES: dict[str, _ThemeColors] = {
	"default": _ThemeColors(
		background=RGBColor(0xFF, 0xFF, 0xFF),
		text=RGBColor(0x1F, 0x29, 0x37),
		muted=RGBColor(0x6B, 0x72, 0x80),
		accent=RGBColor(0x25, 0x63, 0xEB),
		accent_soft=RGBColor(0xDB, 0xEA, 0xFE),
		hero_background=RGBColor(0x1E, 0x3A, 0x8A),
		hero_text=RGBColor(0xFF, 0xFF, 0xFF),
		bullet_background=RGBColor(0xF9, 0xFA, 0xFB),
	),
	"dark": _ThemeColors(
		background=RGBColor(0x0F, 0x17, 0x2A),
		text=RGBColor(0xE2, 0xE8, 0xF0),
		muted=RGBColor(0x94, 0xA3, 0xB8),
		accent=RGBColor(0x60, 0xA5, 0xFA),
		accent_soft=RGBColor(0x1E, 0x29, 0x3B),
		hero_background=RGBColor(0x02, 0x06, 0x17),
		hero_text=RGBColor(0xF8, 0xFA, 0xFC),
		bullet_background=RGBColor(0x1E, 0x29, 0x3B),
	),
}


class SlidePptxService:
	"""Render a `SlideDeckSchema` into a PowerPoint presentation."""

	async def render(
		self,
		deck: SlideDeckSchema,
		theme: str = "default",
	) -> bytes:
		"""Render the deck into a PPTX file.

		WeasyPrint runs synchronously, so the heavy lifting is wrapped
		in `asyncio.to_thread` to keep the event loop responsive.

		Args:
		    deck: The parsed slide deck.
		    theme: The visual theme to apply ("default" or "dark").

		Returns:
		    The PPTX file as bytes.
		"""
		return await asyncio.to_thread(self._render_sync, deck, theme)

	def _render_sync(
		self,
		deck: SlideDeckSchema,
		theme: str,
	) -> bytes:
		"""Synchronously build the PPTX file in memory.

		Args:
		    deck: The parsed slide deck.
		    theme: The visual theme name.

		Returns:
		    The PPTX bytes.
		"""
		presentation: Presentation = Presentation()
		presentation.slide_width = Inches(_SLIDE_WIDTH_INCHES)
		presentation.slide_height = Inches(_SLIDE_HEIGHT_INCHES)
		colors: _ThemeColors = _THEMES.get(theme, _THEMES["default"])

		for slide_data in deck.slides:
			self._add_slide(presentation, slide_data, deck, colors)

		buffer: io.BytesIO = io.BytesIO()
		presentation.save(buffer)
		return buffer.getvalue()

	def _add_slide(
		self,
		presentation: Presentation,
		slide: SlideSchema,
		deck: SlideDeckSchema,
		colors: _ThemeColors,
	) -> None:
		"""Add a single slide to the presentation based on its kind.

		Args:
		    presentation: The PPTX presentation under construction.
		    slide: The schema describing the slide to add.
		    deck: The whole deck (used for footer info).
		    colors: Theme colors to apply.
		"""
		layout = presentation.slide_layouts[_HERO_LAYOUT_INDEX]
		new_slide = presentation.slides.add_slide(layout)
		is_hero: bool = slide.kind == SlideKind.LESSON_TITLE
		bg_color: RGBColor = (
			colors.hero_background if is_hero else colors.background
		)
		self._set_background(new_slide, bg_color)

		if is_hero:
			self._render_hero(new_slide, slide, colors)
		elif slide.kind == SlideKind.BULLETS:
			self._render_bullets(new_slide, slide, colors)
		elif slide.kind == SlideKind.CONCLUSION:
			self._render_content(new_slide, slide, colors, conclusion=True)
		else:
			self._render_content(new_slide, slide, colors, conclusion=False)

		if slide.notes:
			notes_slide = new_slide.notes_slide
			notes_slide.notes_text_frame.text = slide.notes

		if not is_hero:
			self._add_footer(new_slide, deck, slide, colors)

	@staticmethod
	def _set_background(new_slide, color: RGBColor) -> None:
		"""Set a solid background fill on the slide.

		Args:
		    new_slide: The python-pptx slide object.
		    color: The RGB color to apply.
		"""
		fill = new_slide.background.fill
		fill.solid()
		fill.fore_color.rgb = color

	def _render_hero(
		self,
		new_slide,
		slide: SlideSchema,
		colors: _ThemeColors,
	) -> None:
		"""Render a lesson-title (hero) slide.

		Args:
		    new_slide: The python-pptx slide object.
		    slide: The slide schema.
		    colors: Theme colors to apply.
		"""
		if slide.subtitle:
			eyebrow = new_slide.shapes.add_textbox(
				Inches(0.7),
				Inches(2.7),
				Inches(_SLIDE_WIDTH_INCHES - 1.4),
				Inches(0.5),
			)
			self._fill_textbox(
				eyebrow,
				slide.subtitle.upper(),
				font_size=Pt(20),
				color=colors.hero_text,
				bold=False,
				alignment=PP_ALIGN.CENTER,
			)

		title_box = new_slide.shapes.add_textbox(
			Inches(0.7),
			Inches(3.3),
			Inches(_SLIDE_WIDTH_INCHES - 1.4),
			Inches(2.0),
		)
		self._fill_textbox(
			title_box,
			slide.title,
			font_size=Pt(48),
			color=colors.hero_text,
			bold=True,
			alignment=PP_ALIGN.CENTER,
		)

	def _render_content(
		self,
		new_slide,
		slide: SlideSchema,
		colors: _ThemeColors,
		conclusion: bool,
	) -> None:
		"""Render a content (paragraph) slide.

		Args:
		    new_slide: The python-pptx slide object.
		    slide: The slide schema.
		    colors: Theme colors to apply.
		    conclusion: Whether this is the conclusion slide variant.
		"""
		eyebrow_text: str = (
			"CONCLUSÃO"
			if conclusion
			else (
				f"AULA {slide.lesson_index}"
				if slide.lesson_index is not None
				else ""
			)
		)
		if eyebrow_text:
			eyebrow = new_slide.shapes.add_textbox(
				Inches(0.7),
				Inches(0.5),
				Inches(_SLIDE_WIDTH_INCHES - 1.4),
				Inches(0.4),
			)
			self._fill_textbox(
				eyebrow,
				eyebrow_text,
				font_size=Pt(14),
				color=colors.accent,
				bold=True,
				alignment=PP_ALIGN.LEFT,
			)

		title_box = new_slide.shapes.add_textbox(
			Inches(0.7),
			Inches(0.95),
			Inches(_SLIDE_WIDTH_INCHES - 1.4),
			Inches(1.0),
		)
		self._fill_textbox(
			title_box,
			slide.title,
			font_size=Pt(32),
			color=colors.text,
			bold=True,
			alignment=PP_ALIGN.LEFT,
		)

		body_box = new_slide.shapes.add_textbox(
			Inches(0.7),
			Inches(2.1),
			Inches(_SLIDE_WIDTH_INCHES - 1.4),
			Inches(4.4),
		)
		text_frame = body_box.text_frame
		text_frame.word_wrap = True
		text_frame.clear()
		for index, paragraph in enumerate(slide.paragraphs):
			target = (
				text_frame.paragraphs[0]
				if index == 0
				else text_frame.add_paragraph()
			)
			target.alignment = PP_ALIGN.LEFT
			target.space_after = Pt(12)
			run = target.add_run()
			run.text = paragraph
			run.font.size = Pt(20)
			run.font.name = _BODY_FONT
			run.font.color.rgb = colors.text

	def _render_bullets(
		self,
		new_slide,
		slide: SlideSchema,
		colors: _ThemeColors,
	) -> None:
		"""Render a bullet-list slide.

		Args:
		    new_slide: The python-pptx slide object.
		    slide: The slide schema.
		    colors: Theme colors to apply.
		"""
		if slide.lesson_index is not None:
			eyebrow = new_slide.shapes.add_textbox(
				Inches(0.7),
				Inches(0.5),
				Inches(_SLIDE_WIDTH_INCHES - 1.4),
				Inches(0.4),
			)
			self._fill_textbox(
				eyebrow,
				f"AULA {slide.lesson_index}",
				font_size=Pt(14),
				color=colors.accent,
				bold=True,
				alignment=PP_ALIGN.LEFT,
			)

		title_box = new_slide.shapes.add_textbox(
			Inches(0.7),
			Inches(0.95),
			Inches(_SLIDE_WIDTH_INCHES - 1.4),
			Inches(1.0),
		)
		self._fill_textbox(
			title_box,
			slide.title,
			font_size=Pt(32),
			color=colors.text,
			bold=True,
			alignment=PP_ALIGN.LEFT,
		)

		body_box = new_slide.shapes.add_textbox(
			Inches(0.7),
			Inches(2.1),
			Inches(_SLIDE_WIDTH_INCHES - 1.4),
			Inches(4.4),
		)
		text_frame = body_box.text_frame
		text_frame.word_wrap = True
		text_frame.clear()
		for index, bullet in enumerate(slide.bullets):
			target = (
				text_frame.paragraphs[0]
				if index == 0
				else text_frame.add_paragraph()
			)
			target.alignment = PP_ALIGN.LEFT
			target.space_after = Pt(10)
			target.level = 0

			marker_run = target.add_run()
			marker_run.text = "• "
			marker_run.font.size = Pt(20)
			marker_run.font.color.rgb = colors.accent
			marker_run.font.bold = True

			if bullet.term:
				term_run = target.add_run()
				term_run.text = f"{bullet.term}: "
				term_run.font.size = Pt(20)
				term_run.font.bold = True
				term_run.font.color.rgb = colors.accent
				term_run.font.name = _BODY_FONT

			desc_run = target.add_run()
			desc_run.text = bullet.description
			desc_run.font.size = Pt(20)
			desc_run.font.color.rgb = colors.text
			desc_run.font.name = _BODY_FONT

	def _add_footer(
		self,
		new_slide,
		deck: SlideDeckSchema,
		slide: SlideSchema,
		colors: _ThemeColors,
	) -> None:
		"""Add a small footer with the deck title and slide number.

		Args:
		    new_slide: The python-pptx slide object.
		    deck: The full deck.
		    slide: The current slide schema.
		    colors: Theme colors to apply.
		"""
		footer = new_slide.shapes.add_textbox(
			Inches(0.7),
			Inches(_SLIDE_HEIGHT_INCHES - 0.5),
			Inches(_SLIDE_WIDTH_INCHES - 1.4),
			Inches(0.3),
		)
		text_frame = footer.text_frame
		paragraph = text_frame.paragraphs[0]
		paragraph.alignment = PP_ALIGN.LEFT

		deck_run = paragraph.add_run()
		deck_run.text = deck.title
		deck_run.font.size = Pt(10)
		deck_run.font.color.rgb = colors.muted
		deck_run.font.name = _BODY_FONT

		spacer_run = paragraph.add_run()
		spacer_run.text = "    "

		number_run = paragraph.add_run()
		number_run.text = f"{slide.slide_number} / {deck.total}"
		number_run.font.size = Pt(10)
		number_run.font.color.rgb = colors.muted
		number_run.font.name = _BODY_FONT

	@staticmethod
	def _fill_textbox(
		textbox,
		content: str,
		font_size: Pt,
		color: RGBColor,
		bold: bool,
		alignment: PP_ALIGN,
	) -> None:
		"""Populate a textbox with a single styled run.

		Args:
		    textbox: The python-pptx shape.
		    content: The text to display.
		    font_size: The font size to apply.
		    color: The text color.
		    bold: Whether to render the text in bold.
		    alignment: Paragraph alignment to apply.
		"""
		text_frame = textbox.text_frame
		text_frame.word_wrap = True
		text_frame.clear()
		paragraph = text_frame.paragraphs[0]
		paragraph.alignment = alignment
		run = paragraph.add_run()
		run.text = content
		run.font.size = font_size
		run.font.color.rgb = color
		run.font.bold = bold
		run.font.name = _TITLE_FONT
